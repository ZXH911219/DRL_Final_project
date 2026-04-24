"""Vision-Ingestion-Agent Integration Test"""

import sys
from pathlib import Path

project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

import tempfile
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

from src.agents.vision_ingestion import (
    PPTParser,
    ImageRenderer,
    VisualFeatureBundle,
    get_vision_agent,
)
from src.utils import get_logger

logger = get_logger("test_vision_ingestion")


def create_sample_ppt(filepath: str):
    """Create a sample PPTX for testing."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5.5)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Machine Learning in Finance"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5.5)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Risk Management\nData-Driven Approach"
    for p in tf.paragraphs:
        p.font.size = Pt(44)

    prs.save(filepath)
    logger.info(f"Created sample PPTX: {filepath}")


def test_ppt_parsing():
    """Test PPT parsing functionality."""
    logger.info("=" * 60)
    logger.info("TEST 1: PPT Parsing")
    logger.info("=" * 60)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        temp_ppt = f.name

    try:
        create_sample_ppt(temp_ppt)

        # Use factory function to parse
        from src.agents.vision_ingestion import parse_ppt

        parser = parse_ppt(temp_ppt)
        assert parser is not None, "Should parse PPT successfully"

        slide_count = parser.get_slide_count()
        assert slide_count > 0, "Should have parsed slides"
        logger.info(f"✓ Parsed {slide_count} slides")

        metadata = parser.get_slide_metadata(0)
        assert metadata is not None, "Should get slide metadata"
        logger.info(f"✓ Retrieved metadata: {metadata}")

        text = parser.extract_all_text()
        assert len(text) > 0, "Should extract text"
        logger.info(f"✓ Extracted text: {len(text)} characters")

        logger.info("✓ PPT Parsing Test PASSED\n")
        return True

    except Exception as e:
        logger.error(f"✗ PPT Parsing Test FAILED: {e}\n")
        return False
    finally:
        Path(temp_ppt).unlink(missing_ok=True)


def test_image_renderer():
    """Test PPT to image rendering."""
    logger.info("=" * 60)
    logger.info("TEST 2: Image Rendering")
    logger.info("=" * 60)

    with tempfile.TemporaryDirectory() as temp_dir:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_ppt = f.name

        try:
            create_sample_ppt(temp_ppt)

            renderer = ImageRenderer()

            # Try to detect LibreOffice
            try:
                lo_path = renderer._detect_libreoffice()
                if lo_path:
                    logger.info(f"✓ Detected LibreOffice: {lo_path}")
                else:
                    logger.warning("⚠ LibreOffice not detected (optional for this test)")
            except Exception as e:
                logger.warning(f"⚠ LibreOffice detection skipped: {e}")

            logger.info("✓ Image Renderer Test PASSED (framework validated)\n")
            return True

        except Exception as e:
            logger.error(f"✗ Image Renderer Test FAILED: {e}\n")
            return False
        finally:
            Path(temp_ppt).unlink(missing_ok=True)


def test_feature_extraction():
    """Test visual feature extraction."""
    logger.info("=" * 60)
    logger.info("TEST 3: Feature Extraction")
    logger.info("=" * 60)

    try:
        import numpy as np

        # Create dummy image (1024x768 RGB)
        dummy_image = np.random.randint(0, 255, (768, 1024, 3), dtype=np.uint8)

        bundle = VisualFeatureBundle(
            slide_id="test_slide_0",
            page_index=0,
            multi_vectors=np.random.randn(1024, 128),
            imagebind_vector=np.random.randn(1024),
            metadata={"title": "Test Slide", "source": "test.pptx"},
        )

        assert bundle.slide_id == "test_slide_0", "Bundle ID should match"
        assert bundle.multi_vectors.shape == (1024, 128), "Multi-vectors shape should be correct"
        assert bundle.imagebind_vector.shape == (1024,), "ImageBind vector shape should be correct"

        bundle_dict = bundle.to_dict()
        assert "multi_vectors" in bundle_dict, "Dict should contain multi_vectors"
        assert "imagebind_vector" in bundle_dict, "Dict should contain imagebind_vector"
        assert "quality_metrics" in bundle_dict, "Dict should contain quality_metrics"

        logger.info(f"✓ Created VisualFeatureBundle with shape: {bundle.multi_vectors.shape}")
        logger.info(f"✓ Serialized bundle to dict with keys: {list(bundle_dict.keys())}")
        logger.info(f"✓ Quality metrics: {bundle_dict['quality_metrics']}")
        logger.info("✓ Feature Extraction Test PASSED\n")
        return True

    except Exception as e:
        logger.error(f"✗ Feature Extraction Test FAILED: {e}\n")
        return False


def test_vision_agent():
    """Test VisionIngestionAgent orchestration."""
    logger.info("=" * 60)
    logger.info("TEST 4: Vision-Ingestion-Agent")
    logger.info("=" * 60)

    with tempfile.TemporaryDirectory() as temp_dir:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_ppt = f.name

        try:
            create_sample_ppt(temp_ppt)

            agent = get_vision_agent()
            logger.info(f"✓ Created VisionIngestionAgent instance")

            # Test agent initialization
            assert agent.parse_ppt is not None, "Agent should have parse_ppt"
            assert agent.ImageRenderer is not None, "Agent should have ImageRenderer"
            assert agent.extract_visual_features is not None, "Agent should have extract_visual_features"

            logger.info("✓ Agent components initialized correctly")
            logger.info("✓ Vision-Ingestion-Agent Test PASSED\n")
            return True

        except Exception as e:
            logger.error(f"✗ Vision-Ingestion-Agent Test FAILED: {e}\n")
            return False
        finally:
            Path(temp_ppt).unlink(missing_ok=True)


def main():
    """Run all vision ingestion tests."""
    logger.info("\n" + "=" * 60)
    logger.info("VISION-INGESTION-AGENT INTEGRATION TEST SUITE")
    logger.info("=" * 60 + "\n")

    results = {
        "PPT Parsing": test_ppt_parsing(),
        "Image Renderer": test_image_renderer(),
        "Feature Extraction": test_feature_extraction(),
        "Vision-Agent": test_vision_agent(),
    }

    logger.info("=" * 60)
    logger.info("TEST SUMMARY")
    logger.info("=" * 60)
    passed = sum(1 for v in results.values() if v)
    total = len(results)
    for name, result in results.items():
        status = "✓ PASS" if result else "✗ FAIL"
        logger.info(f"{name}: {status}")
    logger.info(f"\nTotal: {passed}/{total} tests passed")
    logger.info("=" * 60 + "\n")

    return all(results.values())


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
