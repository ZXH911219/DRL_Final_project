import os
import time
from pathlib import Path
from typing import List, Dict, Any
import shutil

# DRL Pipeline imports
from src.core.pipeline import get_pipeline, PipelineResult
from src.storage.lancedb_client import LanceDBClient, VectorDocument
from src.models.model_loaders import ColPaliLoader
import numpy as np

def convert_ppt_to_images_windows(ppt_path: str, output_dir: str) -> List[str]:
    """使用 PowerPoint COM 將 PPT 轉為 PDF，再用 PyMuPDF (fitz) 高畫質轉圖片 (避開 Export API 的雷)"""
    print(f"🔄 正在將 PPT 轉為高畫質圖片: {os.path.basename(ppt_path)}")
    try:
        import fitz  # PyMuPDF
        import win32com.client
        import pythoncom
        
        ppt_path_abs = os.path.abspath(ppt_path)
        out_dir_abs = os.path.abspath(output_dir)
        os.makedirs(out_dir_abs, exist_ok=True)
        
        # 1. PPT 轉為 PDF 中介檔 (使用 COM 但不調用任何 UI/Export 屬性)
        pdf_path = os.path.join(out_dir_abs, "temp_render.pdf")
        
        print("   - 步驟 1: 將 PPT 轉換為向量 PDF 中介格式...")
        pythoncom.CoInitialize()
        
        # 啟動 PPT 並設定為完全隱藏
        powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
        # 移除了 Visible 屬性調用，因為有些環境不允許修改
        
        # 開啟 PPT (唯讀, 隱藏, 不開視窗)
        presentation = powerpoint.Presentations.Open(ppt_path_abs, ReadOnly=True, Untitled=False, WithWindow=False)
        
        # 將 PPT 存為 PDF (參數 32 代表 ppSaveAsPDF)
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        
        # 2. 用 PyMuPDF 讀取 PDF 並導出高畫質圖片
        print("   - 步驟 2: 使用 PyMuPDF 將 PDF 渲染為高解析度圖片...")
        doc = fitz.open(pdf_path)
        exported_images = []
        
        # 設定渲染矩陣，放大兩倍 (相當於約 300 DPI)
        zoom_matrix = fitz.Matrix(2.0, 2.0)
        
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=zoom_matrix)
            slide_name = f"slide_{i:03d}.png"
            out_path = os.path.join(out_dir_abs, slide_name)
            pix.save(out_path)
            exported_images.append(out_path)
            
        doc.close()
        
        # 清理暫存的 PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
            
        print(f"✅ 成功渲染了 {len(exported_images)} 張投影片圖片!")
        return exported_images
        
    except Exception as e:
        print(f"❌ 渲染失敗，發生錯誤: {e}")
        return []

def extract_ppt_text(ppt_path: str) -> Dict[int, str]:
    """使用 python-pptx 提取文字"""
    print(f"🔍 正在解析投影片文字結構: {os.path.basename(ppt_path)}")
    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
        slides_text = {}
        
        total_chars = 0
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
            
            slide_text = " ".join(text_parts)
            slides_text[i] = slide_text
            total_chars += len(slide_text)
            
        print(f"✅ 成功解析了 {len(prs.slides)} 頁投影片的文字 (共 {total_chars} 個字元)!")
        return slides_text
        
    except Exception as e:
        print(f"❌ 文字解析失敗: {e}")
        return {}

def test_on_my_ppts():
    """專門為測試用戶 PPT 撰寫的本地測試腳本"""
    ppt_files = ["1141224.pptx", "1150423.pptx"]
    output_base = "data/test_output"
    
    os.makedirs(output_base, exist_ok=True)
    
    # 初始化 LanceDB 連線
    print("⏳ 初始化 LanceDB 向量資料庫...")
    db_client = LanceDBClient(db_path="data/vector_store/lancedb_my_ppts")
    db_client.create_table("ppt_slides")
    
    # 模擬的 ColPali 視覺編碼器 (如果電腦沒裝 Torch / 權重)
    vision_encoder = ColPaliLoader("./models", {"type": "vision"})
    vision_encoder.load()

    total_images_processed = 0
    total_docs_inserted = 0
    
    for ppt_file in ppt_files:
        if not os.path.exists(ppt_file):
            print(f"⚠️ 找不到檔案: {ppt_file}，跳過。")
            continue
            
        print(f"\n" + "="*50)
        print(f"▶️ 開始處理: {ppt_file}")
        print("="*50)
        
        # 1. 萃取純文字內容
        text_dict = extract_ppt_text(ppt_file)
        
        # 2. 轉換為高畫質圖像
        ppt_name = Path(ppt_file).stem
        image_dir = os.path.join(output_base, ppt_name)
        image_paths = convert_ppt_to_images_windows(ppt_file, image_dir)
        
        # 3. 給每頁投影片產生多模態向量並寫入 DB
        if image_paths and text_dict:
            print("🧠 啟動 Vision Agent 提取 1024 區塊的視覺特徵，並寫入 LanceDB...")
            conn = db_client.pool.acquire()
            if not conn:
                print("❌ 拿不到資料庫連線")
                continue
                
            docs = []
            for i, img_path in enumerate(image_paths):
                # 利用 Mock ColPali 提取 (會給出 (1024, 128) 向量)
                multi_vectors = vision_encoder.process_image(img_path)
                
                doc = VectorDocument(
                    doc_id=f"{ppt_name}_slide_{i:03d}",
                    content_type="slide",
                    vectors=multi_vectors, 
                    text_content=text_dict.get(i, ""),
                    metadata={
                        "source": ppt_file,
                        "page": i,
                        "file_name": ppt_name
                    }
                )
                docs.append(doc)
            
            # 寫入 LanceDB
            try:
                db_client.insert_documents(conn, "ppt_slides", docs)
                total_docs_inserted += len(docs)
                print(f"💾 成功將 {len(docs)} 頁投影片索引寫入 LanceDB。")
            finally:
                db_client.pool.release(conn)
                
    if total_docs_inserted > 0:
        # 建立 IVF 索引以加速
        print("\n⚙️ 建立 IVF 向量索引中 (準備快速檢索)...")
        conn = db_client.pool.acquire()
        try:
            db_client.create_ivf_index(conn, "ppt_slides")
            print("✅ 索引建立完成")
            
            # 來試著做一次檢索
            print("\n" + "="*50)
            print("🚀 執行真實的雙階段檢索測試 (Hybrid Search)")
            query_vector = vision_encoder.process_image(image_paths[0]) # 拿第一張圖模擬
            # 對我們擷取的第一頁投影片中的文字下 Query:
            sample_text = text_dict.get(0, "機器學習")[:15] # 拿前 15 個字
            print(f"搜尋 Query 文本: '{sample_text}'")
            
            results = db_client.hybrid_search(
                conn=conn,
                table_name="ppt_slides",
                query_vector=query_vector,
                query_text=sample_text,
                k1=50,
                k2=5,
                fts_weight=0.5,
                vector_weight=0.5
            )
            
            print("\n🏆 Top 5 檢索結果:")
            for i, res in enumerate(results):
                print(f"   [{i+1}] {res.doc_id} (分數: {res.score:.4f}, 來自 {res.stage})")
                print(f"       -> 預覽文字: {res.metadata.get('preview_text', '無')}...")
                
        finally:
            db_client.pool.release(conn)
            
    print("\n🎉 PPT 測試腳本執行完畢！生成的高清圖片保存在 data/test_output/ 中。")

if __name__ == "__main__":
    test_on_my_ppts()
