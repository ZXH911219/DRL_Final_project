"""
PPT File Parser Module
Handle .pptx and .odp format reading and structure extraction.
"""

from pathlib import Path
from typing import Any, Dict, List, Optional

import numpy as np


class PPTParser:
    """Parse PowerPoint files and extract slide information."""

    def __init__(self, file_path: str):
        """
        Initialize PPT parser.

        Args:
            file_path: Path to .pptx or .odp file
        """
        self.file_path = Path(file_path)
        self.format = self._detect_format()
        self._ppt_obj = None

    def _detect_format(self) -> str:
        """Detect PPT file format."""
        suffix = self.file_path.suffix.lower()
        if suffix == ".pptx":
            return "pptx"
        elif suffix == ".odp":
            return "odp"
        else:
            raise ValueError(f"Unsupported format: {suffix}")

    def load(self) -> bool:
        """Load PPT file."""
        try:
            if self.format == "pptx":
                from pptx import Presentation

                self._ppt_obj = Presentation(str(self.file_path))
            elif self.format == "odp":
                # LibreOffice ODP support via uno bridge (optional)
                # For now, raise not implemented
                raise NotImplementedError("ODP support requires LibreOffice UNO bridge")

            return True
        except ImportError:
            print("ERROR: python-pptx not installed. Run: pip install python-pptx")
            return False
        except Exception as e:
            print(f"ERROR loading PPT: {e}")
            return False

    def get_slide_count(self) -> int:
        """Get total number of slides."""
        if not self._ppt_obj:
            return 0
        return len(self._ppt_obj.slides)

    def get_slide_metadata(self, slide_idx: int) -> Dict[str, Any]:
        """
        Get metadata for a specific slide.

        Args:
            slide_idx: 0-based slide index

        Returns:
            Slide metadata dict
        """
        if not self._ppt_obj or slide_idx >= len(self._ppt_obj.slides):
            return {}

        slide = self._ppt_obj.slides[slide_idx]

        metadata = {
            "slide_id": f"slide_{slide_idx}",
            "page_index": slide_idx,
            "slide_layout": slide.slide_layout.name,
            "source_path": str(self.file_path),
            "timestamp": None,
            "shapes_count": len(slide.shapes),
            "has_notes": bool(slide.notes_slide.notes_text_frame.text) if slide.notes_slide else False,
        }

        # Extract text content
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_content.append(shape.text)

        metadata["text_content"] = " ".join(text_content)

        return metadata

    def get_all_slides_metadata(self) -> List[Dict[str, Any]]:
        """Get metadata for all slides."""
        if not self._ppt_obj:
            return []

        return [self.get_slide_metadata(i) for i in range(len(self._ppt_obj.slides))]

    def extract_text_by_slide(self, slide_idx: int) -> str:
        """Extract all text from a slide."""
        if not self._ppt_obj or slide_idx >= len(self._ppt_obj.slides):
            return ""

        slide = self._ppt_obj.slides[slide_idx]
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

        return " ".join(text_parts)

    def extract_all_text(self) -> Dict[int, str]:
        """Extract text from all slides."""
        if not self._ppt_obj:
            return {}

        return {i: self.extract_text_by_slide(i) for i in range(len(self._ppt_obj.slides))}


def parse_ppt(file_path: str) -> Optional[PPTParser]:
    """
    Load and parse a PPT file.

    Args:
        file_path: Path to PPT file

    Returns:
        PPTParser instance or None on error
    """
    parser = PPTParser(file_path)
    if parser.load():
        return parser
    return None
