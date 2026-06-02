"""
Vision-Ingestion-Agent（對齊 openspec/agents.md §1、openspec/specs/specs.md §1.2）。

流程：python-pptx 結構化讀取 → LibreOffice headless 轉 PDF → pdf2image 高 DPI 點陣化
→ ColPali 多向量 (1024, 128) →（可選）ImageBind 占位 → LanceDB 寫入。

* SLA：單頁渲染 < 500ms、特徵提取 < 2s 為目標；超標時於 quality_metrics 標記並記錄 warning。
* patch 順序：與規格一致 **patch_index = gy * 32 + gx**（gy 為列自上而下，gx 為行由左而右），
  `patch_bboxes` 與 `patch_coordinates` 與 `multi_vectors` 列對齊。
"""

from __future__ import annotations

import hashlib
import logging
import os
import re
import shutil
import subprocess
import tempfile
import time
from enum import Enum
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal, Sequence

import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation

from ..storage.lancedb_manager import (
    COLPALI_DIM,
    COLPALI_PATCHES,
    IMAGEBIND_DIM,
    LanceDBManager,
)

logger = logging.getLogger(__name__)

GRID: int = 32
SLA_RENDER_MS: float = 500.0
SLA_ENCODE_MS: float = 2000.0
SLA_IMAGEBIND_MS: float = 300.0


class RenderBackend(str, Enum):
    AUTO = "auto"
    POWERPOINT = "powerpoint"
    LIBREOFFICE = "libreoffice"


def _find_soffice() -> str | None:
    for name in ("soffice", "soffice.exe"):
        p = shutil.which(name)
        if p:
            return p
    win = os.environ.get("ProgramFiles", r"C:\Program Files")
    cand = Path(win) / "LibreOffice" / "program" / "soffice.exe"
    if cand.is_file():
        return str(cand)
    return None


def _ppt_render_backend() -> RenderBackend:
    return RenderBackend.POWERPOINT


def _find_powerpoint_com() -> Any | None:
    """Return a PowerPoint COM application if pywin32 and PowerPoint are available."""
    import pythoncom
    pythoncom.CoInitialize()
    import win32com.client  # type: ignore
    app = win32com.client.DispatchEx("PowerPoint.Application")
    return app


def patch_grid_coordinates_pixel(width: int, height: int) -> list[tuple[int, int]]:
    """
    各 patch 對應 **網格左上角像素座標 (gx, gy)**，索引 i = gy * 32 + gx。
    """
    coords: list[tuple[int, int]] = []
    for gy in range(GRID):
        for gx in range(GRID):
            px = int(round(gx * width / GRID))
            py = int(round(gy * height / GRID))
            coords.append((px, py))
    assert len(coords) == COLPALI_PATCHES
    return coords


def patch_bboxes_normalized(width: int, height: int) -> list[list[float]]:
    """
    Lance `patch_bboxes`：1024 × [x0,y0,x1,y1] **0–1 正規化**，列序 **gy·32+gx**。

    網格均分全圖，與 `specs.md` 證據地圖之 patch→像素映射一致。
    """
    if width <= 0 or height <= 0:
        raise ValueError("寬高必須為正")
    out: list[list[float]] = []
    for gy in range(GRID):
        for gx in range(GRID):
            x0 = gx / GRID
            x1 = (gx + 1) / GRID
            y0 = gy / GRID
            y1 = (gy + 1) / GRID
            out.append([float(x0), float(y0), float(x1), float(y1)])
    return out


def ensure_rgb(image: Image.Image) -> Image.Image:
    if image.mode in ("RGB",):
        return image
    if image.mode == "RGBA":
        bg = Image.new("RGB", image.size, (255, 255, 255))
        bg.paste(image, mask=image.split()[3])
        return bg
    return ImageOps.exif_transpose(image).convert("RGB")


def pptx_to_pdf_libreoffice(pptx_path: Path, out_dir: Path, *, timeout_s: int = 120) -> Path:
    """使用 LibreOffice headless 將 .pptx 轉為同主檔名之 PDF。"""
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError(
            "找不到 LibreOffice (soffice)。請安裝並加入 PATH，或設定正確的安裝路徑。",
        )
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf:impress_pdf_Export",
        "--outdir",
        str(out_dir.resolve()),
        str(pptx_path.resolve()),
    ]
    t0 = time.perf_counter()
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout_s, check=False)
    elapsed_ms = (time.perf_counter() - t0) * 1000
    if proc.returncode != 0:
        raise RuntimeError(
            f"LibreOffice 轉檔失敗 (code={proc.returncode}): {proc.stderr or proc.stdout}",
        )
    pdf_path = out_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.is_file():
        raise FileNotFoundError(f"預期產生 PDF：{pdf_path}（耗時 {elapsed_ms:.1f} ms）")
    logger.info("LibreOffice 轉 PDF 完成：%s（%.1f ms）", pdf_path.name, elapsed_ms)
    return pdf_path


def pptx_to_pdf_powerpoint(pptx_path: Path, out_dir: Path, *, timeout_s: int = 120) -> Path:
    """Use PowerPoint COM to export PPTX to PDF for higher-fidelity rendering.

    This path preserves Office fonts/equations much better than LibreOffice when
    PowerPoint is installed on Windows.
    """
    app = _find_powerpoint_com()
    if app is None:
        raise RuntimeError("PowerPoint COM 不可用")

    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = Path(pptx_path)
    pdf_path = out_dir / f"{pptx_path.stem}.pdf"
    t0 = time.perf_counter()
    presentation = None
    try:
        try:
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
        except Exception:
            pass

        # 1 = msoTrue, 0 = msoFalse
        presentation = app.Presentations.Open(str(pptx_path.resolve()), ReadOnly=True, WithWindow=False)
        # 2 = ppFixedFormatTypePDF
        presentation.ExportAsFixedFormat(str(pdf_path.resolve()), 2)
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass

    elapsed_ms = (time.perf_counter() - t0) * 1000
    if not pdf_path.is_file():
        raise FileNotFoundError(f"PowerPoint 轉檔後找不到 PDF：{pdf_path}")
    logger.info("PowerPoint 轉 PDF 完成：%s（%.1f ms）", pdf_path.name, elapsed_ms)
    return pdf_path


def pptx_to_pdf_best_effort(pptx_path: Path, out_dir: Path, *, timeout_s: int = 120) -> tuple[Path, str]:
    """Export PPTX to PDF using the best available backend."""
    backend = _ppt_render_backend()
    tried: list[str] = []

    if backend in (RenderBackend.AUTO, RenderBackend.POWERPOINT):
        tried.append("powerpoint")
        try:
            return pptx_to_pdf_powerpoint(pptx_path, out_dir, timeout_s=timeout_s), "powerpoint"
        except Exception as exc:
            if backend == RenderBackend.POWERPOINT:
                raise
            logger.warning("PowerPoint 匯出失敗，改用 LibreOffice：%s", exc)

    tried.append("libreoffice")
    pdf_path = pptx_to_pdf_libreoffice(pptx_path, out_dir, timeout_s=timeout_s)
    return pdf_path, "libreoffice"


def pptx_to_png_powerpoint(pptx_path: Path, out_dir: Path, *, timeout_s: int = 120) -> list[Image.Image]:
    """Use PowerPoint COM to export each slide directly to PNG.

    This is the highest-fidelity path when Microsoft PowerPoint is installed.
    """
    app = _find_powerpoint_com()
    if app is None:
        raise RuntimeError("PowerPoint COM 不可用")

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = Path(pptx_path)
    export_dir = out_dir / f"{pptx_path.stem}_slides"
    export_dir.mkdir(parents=True, exist_ok=True)

    t0 = time.perf_counter()
    presentation = None
    try:
        try:
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
        except Exception:
            pass

        try:
            app.Visible = 0
        except Exception:
            try:
                app.Visible = 1
            except Exception:
                logger.debug("無法設定 PowerPoint.Application.Visible 屬性，繼續嘗試打開簡報", exc_info=True)

        try:
            presentation = app.Presentations.Open(str(pptx_path.resolve()), ReadOnly=True, WithWindow=False)
        except Exception:
            # Some PowerPoint installations / sessions disallow WithWindow=False; try with True
            logger.debug("Presentations.Open(..., WithWindow=False) 不可用，改用 WithWindow=True", exc_info=True)
            presentation = app.Presentations.Open(str(pptx_path.resolve()), ReadOnly=True, WithWindow=True)
        slide_count = int(presentation.Slides.Count)
        exported_paths: list[Path] = []
        for idx in range(1, slide_count + 1):
            slide = presentation.Slides.Item(idx)
            try:
                if int(slide.SlideShowTransition.Hidden) != 0:
                    continue
            except Exception:
                pass
            slide_path = export_dir / f"slide_{idx:04d}.png"
            slide.Export(str(slide_path.resolve()), "PNG")
            if not slide_path.is_file():
                raise FileNotFoundError(f"PowerPoint 匯出後找不到 PNG：{slide_path}")
            exported_paths.append(slide_path)

        images: list[Image.Image] = []
        for path in exported_paths:
            with Image.open(path) as img:
                images.append(ensure_rgb(img).copy())
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass

    elapsed_ms = (time.perf_counter() - t0) * 1000
    logger.info("PowerPoint 直接匯出 PNG 完成：%s 頁（%.1f ms）", len(images), elapsed_ms)
    return images


def pdf_to_images_high_fidelity(
    pdf_path: Path,
    *,
    dpi: int = 600,
    fmt: str = "png",
    poppler_path: str | Path | None = None,
) -> list[Image.Image]:
    """以 pdf2image 將 PDF 每頁轉為 PIL Image（RGB）。"""
    from pdf2image import convert_from_path

    resolved_poppler_path: str | None
    if poppler_path is None:
        env_poppler = os.environ.get("POPPLER_PATH")
        resolved_poppler_path = str(Path(env_poppler)) if env_poppler else None
    else:
        resolved_poppler_path = str(Path(poppler_path))

    t0 = time.perf_counter()
    images = convert_from_path(
        str(pdf_path),
        dpi=dpi,
        fmt=fmt,
        thread_count=1,
        poppler_path=resolved_poppler_path,
    )
    elapsed_ms = (time.perf_counter() - t0) * 1000
    logger.info("pdf2image：%s 頁，DPI=%s，%.1f ms", len(images), dpi, elapsed_ms)
    return [ensure_rgb(im) for im in images]


def _close_images(images: Sequence[Image.Image]) -> None:
    for image in images:
        try:
            image.close()
        except Exception:
            pass


@dataclass
class SlideMetadata:
    page_index: int
    slide_layout_name: str | None
    has_notes: bool
    slide_title: str | None = None
    body_text: str | None = None
    notes_text: str | None = None


def read_pptx_slide_metadata(pptx_path: Path) -> list[SlideMetadata]:
    prs = Presentation(str(pptx_path))
    meta: list[SlideMetadata] = []
    for i, slide in enumerate(prs.slides):
        try:
            sldId = prs.slides._sldIdLst[i]
            if sldId.get("show") == "0":
                continue
        except Exception:
            pass
            
        layout = slide.slide_layout.name if slide.slide_layout is not None else None
        notes = slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None
        slide_title, body_text, notes_text = extract_slide_text_sections(slide)
        meta.append(
            SlideMetadata(
                page_index=i,
                slide_layout_name=layout,
                has_notes=notes,
                slide_title=slide_title,
                body_text=body_text,
                notes_text=notes_text,
            )
        )
    return meta


def _normalize_text(text: str | None) -> str:
    if not text:
        return ""
    return re.sub(r"\s+", " ", str(text)).strip()


def _shape_text_lines(shape: Any) -> list[str]:
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
        for paragraph in shape.text_frame.paragraphs:
            text = _normalize_text(paragraph.text)
            if text:
                lines.append(text)
    if getattr(shape, "has_table", False) and shape.table is not None:
        for row in shape.table.rows:
            for cell in row.cells:
                cell_text = _normalize_text(cell.text)
                if cell_text:
                    lines.append(cell_text)
    return lines


def extract_slide_text_sections(slide: Any) -> tuple[str | None, str | None, str | None]:
    title_text = ""
    if getattr(slide.shapes, "title", None) is not None:
        title_text = _normalize_text(slide.shapes.title.text)

    body_lines: list[str] = []
    for shape in slide.shapes:
        if getattr(slide.shapes, "title", None) is not None and shape == slide.shapes.title:
            continue
        for line in _shape_text_lines(shape):
            if line:
                body_lines.append(line)

    notes_text = ""
    if getattr(slide, "has_notes_slide", False) and slide.notes_slide is not None:
        notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
        if notes_frame is not None:
            notes_parts = [_normalize_text(paragraph.text) for paragraph in notes_frame.paragraphs]
            notes_text = " ".join(part for part in notes_parts if part)

    body_text = " ".join(body_lines)
    return (
        title_text or None,
        body_text or None,
        notes_text or None,
    )


def build_slide_fts_text(title_text: str | None, body_text: str | None, notes_text: str | None) -> str | None:
    parts: list[str] = []
    if title_text:
        parts.append(f"標題: {_normalize_text(title_text)}")
    if body_text:
        parts.append(f"內文: {_normalize_text(body_text)}")
    if notes_text:
        parts.append(f"備註: {_normalize_text(notes_text)}")
    combined = "\n".join(parts).strip()
    return combined or None


@dataclass
class VisualFeatureBundle:
    """agents.md §1.4 輸出結構（記憶體側；寫入 Lance 時另轉欄位）。"""

    slide_id: str
    page_index: int
    multi_vectors: np.ndarray
    patch_coordinates: list[tuple[int, int]]
    imagebind_vector: np.ndarray | None
    metadata: dict[str, Any] = field(default_factory=dict)
    quality_metrics: dict[str, float] = field(default_factory=dict)


class ColPaliEncoder(ABC):
    """ColPali 視覺編碼器抽象；輸出 float32 (1024, 128)。"""

    @abstractmethod
    def encode(self, image: Image.Image) -> tuple[np.ndarray, dict[str, Any]]:
        """回傳 (multi_vectors, encoder_info)。"""
        raise NotImplementedError


class StubColPaliEncoder(ColPaliEncoder):
    """可重現假資料，供無 GPU／無權重時打通管線。"""

    def __init__(self, seed: int = 42) -> None:
        self._rng = np.random.default_rng(seed)

    def encode(self, image: Image.Image) -> tuple[np.ndarray, dict[str, Any]]:
        _ = image
        v = self._rng.standard_normal((COLPALI_PATCHES, COLPALI_DIM), dtype=np.float32)
        v /= np.maximum(np.linalg.norm(v, axis=1, keepdims=True), 1e-8)
        return v, {"backend": "stub"}


class HFColPaliEncoder(ColPaliEncoder):
    """
    Hugging Face `ColPaliForRetrieval`（需 `transformers` 支援 ColPali、以及 torch）。

    若模型輸出之 token 數或隱藏維度非 1024×128，會以抽樣／截斷方式對齊專案固定網格（並記錄 warning）。
    """

    def __init__(
        self,
        model_id: str = "vidore/colpali-v1.3-hf",
        *,
        device: str | None = None,
        dtype: Literal["auto", "float16", "bfloat16", "float32"] = "auto",
    ) -> None:
        import torch
        from transformers import ColPaliForRetrieval, ColPaliProcessor

        self._torch = torch
        if device is None:
            device = "cuda" if torch.cuda.is_available() else "cpu"
        self._device = device

        torch_dtype: Any = None
        if dtype == "float16":
            torch_dtype = torch.float16
        elif dtype == "bfloat16":
            torch_dtype = torch.bfloat16
        elif dtype == "float32":
            torch_dtype = torch.float32
        elif dtype == "auto":
            torch_dtype = "auto"

        self._processor = ColPaliProcessor.from_pretrained(model_id)
        if device == "cpu":
            kw: dict[str, Any] = {}
            if isinstance(torch_dtype, str):
                kw["torch_dtype"] = torch.float32
            elif torch_dtype is not None:
                kw["torch_dtype"] = torch_dtype
            self._model = ColPaliForRetrieval.from_pretrained(model_id, **kw)
            self._model.to("cpu")
        else:
            kw = {"device_map": "auto"}
            if torch_dtype is not None:
                kw["torch_dtype"] = torch_dtype
            self._model = ColPaliForRetrieval.from_pretrained(model_id, **kw)
        self._model.eval()
        self._model_id = model_id

    def encode(self, image: Image.Image) -> tuple[np.ndarray, dict[str, Any]]:
        torch = self._torch
        rgb = ensure_rgb(image)
        inputs = self._processor(images=[rgb], return_tensors="pt")
        device = next(self._model.parameters()).device
        inputs = {k: v.to(device) for k, v in inputs.items()}
        t0 = time.perf_counter()
        with torch.inference_mode():
            out = self._model(**inputs)
            emb_tensor = getattr(out, "embeddings", None)
            if emb_tensor is None:
                raise RuntimeError(
                    "ColPali 模型輸出缺少 `embeddings` 欄位；請升級 transformers 或檢查模型類別。",
                )
            emb = emb_tensor[0].detach().float().cpu().numpy()
        elapsed_ms = (time.perf_counter() - t0) * 1000
        aligned, info = _align_colpali_tokens_to_grid(emb)
        info.update({"backend": "hf_colpali", "model_id": self._model_id, "encode_ms": elapsed_ms})
        return aligned, info


def _align_colpali_tokens_to_grid(emb: np.ndarray) -> tuple[np.ndarray, dict[str, Any]]:
    """
    emb: (L, H) float；對齊為 (1024, 128)。
    """
    if emb.ndim != 2:
        raise ValueError(f"預期 2D embeddings，收到 shape={emb.shape}")
    l, h = emb.shape
    info: dict[str, Any] = {"source_tokens": int(l), "source_hidden": int(h)}
    e = np.ascontiguousarray(emb, dtype=np.float32)
    if h >= COLPALI_DIM:
        if h != COLPALI_DIM:
            logger.warning("ColPali hidden=%s 與專案 COLPALI_DIM=%s 不符，截斷前 %s 維", h, COLPALI_DIM, COLPALI_DIM)
        e = e[:, :COLPALI_DIM]
    else:
        pad = np.zeros((l, COLPALI_DIM - h), dtype=np.float32)
        e = np.concatenate([e, pad], axis=1)
        logger.warning("ColPali hidden=%s 小於 %s，已零填充", h, COLPALI_DIM)

    if l == COLPALI_PATCHES:
        return e, info

    idx = np.linspace(0, l - 1, COLPALI_PATCHES).astype(np.int64)
    resampled = e[idx]
    logger.warning(
        "ColPali token 數=%s 與專案 1024 不符，已沿 token 軸線性重採樣對齊 32×32 語意網格",
        l,
    )
    info["resampled"] = True
    return resampled, info


class ImageBindEncoderPlaceholder:
    """
    ImageBind 對齊占位（specs §2.3.1：與 ColPali 分路）；啟用時回傳零向量並標記 placeholder。
    """

    def __init__(self, enabled: bool = False) -> None:
        self.enabled = enabled

    def encode_image(self, image: Image.Image) -> tuple[np.ndarray | None, dict[str, Any]]:
        if not self.enabled:
            return None, {"imagebind": "skipped"}
        t0 = time.perf_counter()
        _ = ensure_rgb(image)
        elapsed_ms = (time.perf_counter() - t0) * 1000
        vec = np.zeros((IMAGEBIND_DIM,), dtype=np.float32)
        return vec, {
            "imagebind": "placeholder",
            "encode_ms": elapsed_ms,
            "note": "請替換為實際 ImageBind / 影像塔編碼器；與 ColPali 不得混用於 MaxSim。",
        }


def default_slide_id(deck_id: str, page_index: int, source_path: Path) -> str:
    h = hashlib.sha256(f"{source_path.resolve()}:{page_index}".encode("utf-8")).hexdigest()[:12]
    safe_deck = deck_id.replace("|", "_")
    return f"{safe_deck}|p{page_index:04d}|{h}"


def _safe_artifact_name(name: str) -> str:
    """將 slide_id 轉成可在 Windows / POSIX 安全使用的檔名。"""
    return re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", name).strip("._ ") or "slide"


def slide_image_path(deck_id: str, slide_id: str, root: Path | None = None) -> Path:
    base_dir = Path(root) if root is not None else Path("artifacts")
    return base_dir / "slide_images" / _safe_artifact_name(deck_id) / f"{_safe_artifact_name(slide_id)}.png"


class VisionIngestionAgent:
    """
    Vision-Ingestion-Agent：渲染 + ColPali +（可選）ImageBind + LanceDB。
    """

    def __init__(
        self,
        lance: LanceDBManager,
        *,
        colpali: ColPaliEncoder | None = None,
        imagebind: ImageBindEncoderPlaceholder | None = None,
        dpi: int = 600,
        poppler_path: str | Path | None = None,
    ) -> None:
        self._lance = lance
        self._colpali = colpali or StubColPaliEncoder()
        self._imagebind = imagebind or ImageBindEncoderPlaceholder(enabled=False)
        self._dpi = dpi
        self._poppler_path = poppler_path

    @property
    def lance(self) -> LanceDBManager:
        return self._lance

    def render_pptx_to_images(
        self,
        pptx_path: Path,
        *,
        workdir: Path | None = None,
    ) -> tuple[list[Image.Image], Path | None, str]:
        """
        PowerPoint 直接輸出 PNG；若失敗則退回 LibreOffice → PDF → pdf2image。

        回傳 (每頁 RGB PIL, 產生之 PDF 路徑或 None, render backend)。
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.is_file():
            raise FileNotFoundError(pptx_path)

        if workdir is None:
            workdir = Path(tempfile.mkdtemp(prefix="vision_ingest_"))
        else:
            workdir = Path(workdir)
            workdir.mkdir(parents=True, exist_ok=True)

        backend = _ppt_render_backend()
        if backend in (RenderBackend.AUTO, RenderBackend.POWERPOINT):
            try:
                images = pptx_to_png_powerpoint(pptx_path, workdir)
                return images, None, "powerpoint"
            except Exception as exc:
                if backend == RenderBackend.POWERPOINT:
                    raise
                logger.warning("PowerPoint 直接匯出 PNG 失敗，改用 LibreOffice：%s", exc)

        pdf, pdf_backend = pptx_to_pdf_best_effort(pptx_path, workdir)
        images = pdf_to_images_high_fidelity(
            pdf,
            dpi=self._dpi,
            poppler_path=self._poppler_path,
        )
        return images, pdf, pdf_backend

    def ingest_images(
        self,
        *,
        deck_id: str,
        source_path: Path,
        images: Sequence[Image.Image],
        slide_metadata: Sequence[SlideMetadata] | None = None,
        page_offset: int = 0,
        artifact_root: Path | None = None,
        progress_callback: Any = None,
    ) -> list[VisualFeatureBundle]:
        """
        已具備渲染影像時，直接編碼並寫入 Lance（每張一列）。
        """
        source_path = Path(source_path)
        metas = slide_metadata
        if metas is not None and len(metas) != len(images):
            raise ValueError("slide_metadata 長度必須與 images 相同")

        bundles: list[VisualFeatureBundle] = []
        lance_rows: list[dict[str, Any]] = []

        for i, img in enumerate(images):
            if progress_callback:
                progress_callback(i, len(images), f"正在分析第 {i+1} / {len(images)} 頁 (AI 視覺編碼中)...")
            page_index = page_offset + i
            w, h = img.size
            patch_coords = patch_grid_coordinates_pixel(w, h)
            bboxes_norm = patch_bboxes_normalized(w, h)

            t_r0 = time.perf_counter()
            rgb = ensure_rgb(img)
            render_ms = (time.perf_counter() - t_r0) * 1000
            if render_ms > SLA_RENDER_MS:
                logger.warning(
                    "單頁 RGB 正規化耗時 %.1f ms，超過 SLA 渲染目標 %.0f ms（完整轉檔請見日誌）",
                    render_ms,
                    SLA_RENDER_MS,
                )

            t_e0 = time.perf_counter()
            multi, enc_info = self._colpali.encode(rgb)
            encode_ms = (time.perf_counter() - t_e0) * 1000
            if encode_ms > SLA_ENCODE_MS:
                logger.warning(
                    "ColPali 編碼 %.1f ms 超過 SLA 目標 %.0f ms",
                    encode_ms,
                    SLA_ENCODE_MS,
                )

            ib_vec, ib_info = self._imagebind.encode_image(rgb)
            ib_ms = float(ib_info.get("encode_ms", 0.0))
            if self._imagebind.enabled and ib_ms > SLA_IMAGEBIND_MS:
                logger.warning(
                    "ImageBind 占位 %.1f ms 超過參考 SLA %.0f ms",
                    ib_ms,
                    SLA_IMAGEBIND_MS,
                )

            slide_id = default_slide_id(deck_id, page_index, source_path)
            slide_path = slide_image_path(deck_id, slide_id, artifact_root)
            slide_path.parent.mkdir(parents=True, exist_ok=True)
            ensure_rgb(img).save(slide_path, format="PNG")
            layout_name = None
            has_notes = False
            slide_title = None
            body_text = None
            notes_text = None
            if metas is not None:
                layout_name = metas[i].slide_layout_name
                has_notes = metas[i].has_notes
                slide_title = metas[i].slide_title
                body_text = metas[i].body_text
                notes_text = metas[i].notes_text

            fts_text = build_slide_fts_text(slide_title, body_text, notes_text)

            coverage = 1.0
            qm: dict[str, float] = {
                "render_normalize_ms": float(render_ms),
                "colpali_encode_ms": float(encode_ms),
                "vector_coverage_ratio": float(coverage),
                "slide_width_px": float(w),
                "slide_height_px": float(h),
                "dpi": float(self._dpi),
            }
            if ib_ms:
                qm["imagebind_encode_ms"] = float(ib_ms)

            meta = {
                "deck_id": deck_id,
                "source_path": str(source_path.resolve()),
                "slide_layout": layout_name,
                "has_notes": has_notes,
                "slide_title": slide_title,
                "body_text": body_text,
                "notes_text": notes_text,
                "encoder": enc_info,
                "imagebind": ib_info,
                "slide_image_path": str(slide_path.resolve()),
            }

            bundle = VisualFeatureBundle(
                slide_id=slide_id,
                page_index=page_index,
                multi_vectors=multi,
                patch_coordinates=patch_coords,
                imagebind_vector=ib_vec,
                metadata=meta,
                quality_metrics=qm,
            )
            bundles.append(bundle)

            row: dict[str, Any] = {
                "slide_id": slide_id,
                "deck_id": deck_id,
                "page_index": int(page_index),
                "source_path": str(source_path.resolve()),
                "colpali_multi": multi,
                "patch_bboxes": bboxes_norm,
                "fts_text": fts_text,
                "quality_metrics": qm,
            }
            if ib_vec is not None:
                row["imagebind_vec"] = ib_vec
            lance_rows.append(row)

        self._lance.add(lance_rows, mode="append")
        return bundles

    def ingest_pptx(
        self,
        pptx_path: Path,
        deck_id: str,
        *,
        workdir: Path | None = None,
        keep_artifacts: bool = False,
        artifact_root: Path | None = None,
        progress_callback: Any = None,
    ) -> list[VisualFeatureBundle]:
        """
        端到端：pptx 結構讀取 + 高畫質渲染 + 編碼 + LanceDB。
        """
        pptx_path = Path(pptx_path)
        
        if progress_callback:
            progress_callback(0, 1, "正在從投影片擷取結構化文字...")
            
        meta = read_pptx_slide_metadata(pptx_path)

        tmp: Path
        created_tmp = False
        if workdir is None:
            tmp = Path(tempfile.mkdtemp(prefix="vision_pptx_"))
            created_tmp = True
        else:
            tmp = Path(workdir)
            tmp.mkdir(parents=True, exist_ok=True)

        try:
            if progress_callback:
                progress_callback(0, 1, "正在啟動 PowerPoint 渲染高畫質圖片 (可能需要一兩分鐘)...")
            t0 = time.perf_counter()
            images, pdf_path, render_backend = self.render_pptx_to_images(pptx_path, workdir=tmp)
            render_total_ms = (time.perf_counter() - t0) * 1000
            per_page_render_ms = render_total_ms / max(len(images), 1)
            if per_page_render_ms > SLA_RENDER_MS:
                logger.warning(
                    "平均每頁渲染鏈路 %.1f ms（backend=%s, DPI=%s），高於 SLA 目標 %.0f ms",
                    per_page_render_ms,
                    render_backend,
                    self._dpi,
                    SLA_RENDER_MS,
                )

            meta_list = list(meta)
            if len(meta_list) != len(images):
                logger.warning(
                    "python-pptx 投影片數=%s 與 PDF 頁數=%s 不一致，將擴充／截斷 metadata 對齊",
                    len(meta_list),
                    len(images),
                )
            while len(meta_list) < len(images):
                j = len(meta_list)
                meta_list.append(SlideMetadata(page_index=j, slide_layout_name=None, has_notes=False))
            if len(meta_list) > len(images):
                meta_list = meta_list[: len(images)]

            bundles = self.ingest_images(
                deck_id=deck_id,
                source_path=pptx_path,
                images=list(images),
                slide_metadata=meta_list,
                page_offset=0,
                artifact_root=artifact_root,
                progress_callback=progress_callback,
            )
            for b in bundles:
                if pdf_path is not None:
                    b.metadata["render_pdf_path"] = str(pdf_path)
                b.metadata["render_backend"] = render_backend
                b.metadata["render_total_ms"] = render_total_ms
            return bundles
        finally:
            if created_tmp and not keep_artifacts and tmp.is_dir():
                shutil.rmtree(tmp, ignore_errors=True)

    def ingest_pdf(
        self,
        pdf_path: Path,
        deck_id: str,
        *,
        dpi: int | None = None,
        workdir: Path | None = None,
        artifact_root: Path | None = None,
        progress_callback: Any = None,
    ) -> list[VisualFeatureBundle]:
        """
        端到端：PDF 轉圖片（高保真） + 編碼 + LanceDB。
        
        使用指定 DPI（預設為 self._dpi）從 PDF 每頁生成高解析度影像，
        然後進行 ColPali 編碼並寫入 Lance。
        """
        pdf_path = Path(pdf_path)
        if not pdf_path.is_file():
            raise FileNotFoundError(f"PDF 檔案不存在：{pdf_path}")
        
        if dpi is None:
            dpi = self._dpi

        try:
            t0 = time.perf_counter()
            images = pdf_to_images_high_fidelity(pdf_path, dpi=dpi)
            render_total_ms = (time.perf_counter() - t0) * 1000
            per_page_render_ms = render_total_ms / max(len(images), 1)
            
            if per_page_render_ms > SLA_RENDER_MS:
                logger.warning(
                    "平均每頁 PDF 轉圖片 %.1f ms（DPI=%s），高於 SLA 目標 %.0f ms",
                    per_page_render_ms,
                    dpi,
                    SLA_RENDER_MS,
                )

            # 嘗試從每頁影像做 OCR（優先 pytesseract，否則 fallback 到 easyocr）
            ocr_texts: list[str | None] = [None] * len(images)
            ocr_available = False
            try:
                import pytesseract  # type: ignore

                def _do_ocr_pyt(img: Image.Image) -> str:
                    try:
                        return pytesseract.image_to_string(img, lang=(os.environ.get("TESSERACT_LANG") or "chi_sim+eng"))
                    except Exception:
                        return pytesseract.image_to_string(img)

                ocr_available = True
                for i, img in enumerate(images):
                    try:
                        ocr_texts[i] = _do_ocr_pyt(img.convert("RGB"))
                    except Exception:
                        ocr_texts[i] = None
            except Exception:
                # pytesseract not available or failed; try easyocr
                try:
                    import easyocr  # type: ignore

                    reader = easyocr.Reader(["ch_sim", "en"], gpu=False)
                    ocr_available = True
                    for i, img in enumerate(images):
                        try:
                            res = reader.readtext(np.asarray(img))
                            texts = [t[1] for t in res if t and len(t) >= 1]
                            ocr_texts[i] = "\n".join(texts) if texts else None
                        except Exception:
                            ocr_texts[i] = None
                except Exception:
                    ocr_available = False

            meta_list: list[SlideMetadata] = []
            for i in range(len(images)):
                page_ocr = (ocr_texts[i] or "").strip() if ocr_available else ""
                lines = [ln.strip() for ln in page_ocr.splitlines() if ln.strip()]
                title = None
                body = None
                if lines:
                    # Heuristic: first non-empty OCR line is treated as title
                    title = lines[0][:200]
                    if len(lines) > 1:
                        body = " ".join(lines[1:])[:2000]

                meta_list.append(
                    SlideMetadata(
                        page_index=i,
                        slide_layout_name=None,
                        has_notes=False,
                        slide_title=title,
                        body_text=body,
                        notes_text=None,
                    )
                )

            bundles = self.ingest_images(
                deck_id=deck_id,
                source_path=pdf_path,
                images=list(images),
                slide_metadata=meta_list,
                page_offset=0,
                artifact_root=artifact_root,
                progress_callback=progress_callback,
            )
            
            for b in bundles:
                b.metadata["source_format"] = "pdf"
                b.metadata["pdf_path"] = str(pdf_path)
                b.metadata["render_total_ms"] = render_total_ms
                b.metadata["render_dpi"] = dpi
            
            return bundles
        except Exception as exc:
            logger.error("PDF 匯入失敗：%s", exc, exc_info=True)
            raise


def build_colpali_encoder_from_env() -> ColPaliEncoder:
    """COLPALI_BACKEND=stub|hf；HF 時可設 COLPALI_MODEL、COLPALI_DEVICE、COLPALI_DTYPE。"""
    backend = os.environ.get("COLPALI_BACKEND", "stub").strip().lower()
    if backend in ("stub", "fake", "dummy"):
        seed = int(os.environ.get("COLPALI_STUB_SEED", "42"))
        return StubColPaliEncoder(seed=seed)
    if backend in ("hf", "huggingface", "transformers"):
        model = os.environ.get("COLPALI_MODEL", "vidore/colpali-v1.3-hf")
        device = os.environ.get("COLPALI_DEVICE")
        dtype = os.environ.get("COLPALI_DTYPE", "auto")
        try:
            return HFColPaliEncoder(model, device=device, dtype=dtype)  # type: ignore[arg-type]
        except ImportError as e:
            raise ImportError(
                "COLPALI_BACKEND=hf 需要安裝 torch 與含 ColPali 的 transformers "
                "（建議：`pip install -e \".[vision-ml]\"`）",
            ) from e
    raise ValueError(f"未知 COLPALI_BACKEND={backend!r}（stub 或 hf）")
